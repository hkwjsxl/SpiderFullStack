from docx import Document
from pptx import Presentation


# docx_obj = Document('./练习1.docx')
# for para in docx_obj.paragraphs:
#     for run in para.runs:
#         print(run.text)
#
# docx_obj = Document('./练习2.docx')
# for para in docx_obj.paragraphs:
#     for run in para.runs:
#         print(run.text)
# table = docx_obj.tables[0]
# for row in table.rows:
#     for cell in row.cells:
#         print(cell.text)
# print(''.center(50, '-'))
# print(table.rows[1].cells[2].text)


# pptx_obj = Presentation('./数据报告模板.pptx')
# slide = pptx_obj.slides[1]
# for shape in slide.shapes:
#     if shape.has_text_frame:
#         for para in shape.text_frame.paragraphs:
#             for run in para.runs:
#                 if run.text.endswith('情况'):
#                     run.text = run.text.replace('情况', '分析')
# pptx_obj.save('数据报告模板_new.pptx')
#
#

def create_pptx(name):
    pptx_obj = Presentation('./职业奖励.pptx')
    slide = pptx_obj.slides[0]
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.endswith('xxx'):
                        run.text = name
    pptx_obj.save(f'{name}.pptx')


docx_obj = Document('./练习2.docx')
table = docx_obj.tables[0]
for row in table.rows:
    for cell in row.cells:
        create_pptx(cell.text)